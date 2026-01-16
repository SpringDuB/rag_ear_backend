from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import csv
from typing import Iterable, List


class DocumentParseError(RuntimeError):
    pass


@dataclass
class DocumentParser:
    ocr_lang: str = "chi_sim+eng"
    ocr_min_chars: int = 20

    def parse(self, path: Path | str) -> str:
        file_path = Path(path)
        if not file_path.exists():
            raise DocumentParseError(f"File not found: {file_path}")

        ext = file_path.suffix.lower()
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        if ext == ".docx":
            return self._parse_docx(file_path)
        if ext == ".csv":
            return self._parse_csv(file_path)
        if ext == ".xlsx":
            return self._parse_xlsx(file_path)
        if ext == ".pptx":
            return self._parse_pptx(file_path)

        raise DocumentParseError(f"Unsupported file type: {ext}")

    def _parse_pdf(self, file_path: Path) -> str:
        try:
            import pdfplumber
        except ImportError as exc:
            raise DocumentParseError("Missing dependency: pdfplumber") from exc

        parts: List[str] = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                page_parts: List[str] = []
                text = page.extract_text() or ""
                if text.strip():
                    page_parts.append(text.strip())

                tables = page.extract_tables() or []
                for table in tables:
                    if table:
                        page_parts.append(self._format_table(table))

                if page_parts:
                    parts.append(f"[Page {i}]")
                    parts.extend(page_parts)

        joined = "\n\n".join(parts).strip()
        if len(joined) >= self.ocr_min_chars:
            return joined

        ocr_text = self._ocr_pdf(file_path)
        return ocr_text if ocr_text.strip() else joined

    def _ocr_pdf(self, file_path: Path) -> str:
        try:
            from pdf2image import convert_from_path
        except ImportError as exc:
            raise DocumentParseError("Missing dependency: pdf2image") from exc

        try:
            import pytesseract
        except ImportError as exc:
            raise DocumentParseError("Missing dependency: pytesseract") from exc

        pages = convert_from_path(str(file_path))
        parts: List[str] = []
        for i, image in enumerate(pages, start=1):
            text = pytesseract.image_to_string(image, lang=self.ocr_lang) or ""
            if text.strip():
                parts.append(f"[Page {i} OCR]")
                parts.append(text.strip())
        return "\n\n".join(parts).strip()

    def _parse_docx(self, file_path: Path) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise DocumentParseError("Missing dependency: python-docx") from exc

        doc = Document(str(file_path))
        parts: List[str] = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            parts.append(self._format_table([[cell.text for cell in row.cells] for row in table.rows]))

        return "\n\n".join(parts).strip()

    def _parse_csv(self, file_path: Path) -> str:
        rows: List[List[str]] = []
        with file_path.open("r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append([cell if cell is not None else "" for cell in row])
        return self._format_table(rows)

    def _parse_xlsx(self, file_path: Path) -> str:
        try:
            import openpyxl
        except ImportError as exc:
            raise DocumentParseError("Missing dependency: openpyxl") from exc

        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        parts: List[str] = []
        for sheet in wb.worksheets:
            rows: List[List[str]] = []
            for row in sheet.iter_rows(values_only=True):
                rows.append([self._cell_to_text(cell) for cell in row])
            if rows:
                parts.append(f"[Sheet {sheet.title}]")
                parts.append(self._format_table(rows))
        return "\n\n".join(parts).strip()

    def _parse_pptx(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DocumentParseError("Missing dependency: python-pptx") from exc

        pres = Presentation(str(file_path))
        parts: List[str] = []
        for i, slide in enumerate(pres.slides, start=1):
            slide_parts: List[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        slide_parts.append(text)
                if getattr(shape, "has_table", False):
                    table = shape.table
                    rows = []
                    for r in range(len(table.rows)):
                        row = []
                        for c in range(len(table.columns)):
                            cell = table.cell(r, c)
                            row.append(cell.text if cell.text is not None else "")
                        rows.append(row)
                    if rows:
                        slide_parts.append(self._format_table(rows))
            if slide_parts:
                parts.append(f"[Slide {i}]")
                parts.extend(slide_parts)
        return "\n\n".join(parts).strip()

    def _format_table(self, rows: Iterable[Iterable[str]]) -> str:
        lines = []
        for row in rows:
            safe_row = [(cell or "").strip() for cell in row]
            lines.append("\t".join(safe_row))
        return "\n".join(lines).strip()

    def _cell_to_text(self, value: object) -> str:
        if value is None:
            return ""
        return str(value)
